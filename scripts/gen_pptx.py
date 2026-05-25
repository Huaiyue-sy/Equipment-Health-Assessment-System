#!/usr/bin/env python3
"""Generate iHealthSim presentation PPTX from architecture documentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
C_DARK   = RGBColor(0x0F, 0x17, 0x2A)
C_BRAND  = RGBColor(0x1A, 0x73, 0xE8)
C_BRAND2 = RGBColor(0x4F, 0x46, 0xE5)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BODY   = RGBColor(0x33, 0x41, 0x55)
C_MUTED  = RGBColor(0x64, 0x74, 0x8B)
C_BG     = RGBColor(0xF8, 0xFA, 0xFC)
C_GREEN  = RGBColor(0x10, 0xB9, 0x81)
C_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
C_RED    = RGBColor(0xEF, 0x44, 0x44)
C_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
C_LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)
C_LIGHT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
C_LIGHT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)
C_LIGHT_RED   = RGBColor(0xFE, 0xE2, 0xE2)
C_LIGHT_PURPLE= RGBColor(0xED, 0xE9, 0xFE)
C_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
C_CODE_BG = RGBColor(0x1E, 0x29, 0x3B)
C_CODE_FG = RGBColor(0xE2, 0xE8, 0xF0)

# ── Slide dimensions (16:9) ──
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── Helper Functions ──

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=C_BODY, bold=False, align=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    if font_name:
        p.font.name = font_name
    return txBox

def add_rich_text(slide, left, top, width, height, segments, align=PP_ALIGN.LEFT):
    """segments: list of (text, size, color, bold)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for i, seg in enumerate(segments):
        text, size, color, bold = seg
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_rect_text(slide, left, top, width, height, text, fill_color=C_WHITE,
                  border_color=C_BORDER, font_size=12, font_color=C_BODY, bold=False, align=PP_ALIGN.CENTER):
    shape = add_rect(slide, left, top, width, height, fill_color, border_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    # vertical center
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.auto_size = None
    return shape

def add_card(slide, left, top, width, height, title, body_lines,
             icon="", icon_color=C_BRAND, title_size=14, body_size=11):
    """Add a card with title and bullet body."""
    shape = add_rect(slide, left, top, width, height, C_WHITE, C_BORDER, Pt(1))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.1)

    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.color.rgb = C_DARK
    p.font.bold = True
    p.space_after = Pt(6)

    # Body
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(body_size)
        p.font.color.rgb = C_MUTED
        p.space_after = Pt(2)
        p.level = 0
    return shape

def add_code_block(slide, left, top, width, height, code_text):
    """Add a code block with dark background."""
    shape = add_rect(slide, left, top, width, height, C_CODE_BG, None)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(10)
    p.font.color.rgb = C_CODE_FG
    p.font.name = "Courier New"
    p.space_after = Pt(1)
    return shape

def add_table(slide, left, top, col_widths, headers, rows):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(sum(col_widths)), Inches(0.35 * n_rows))
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Inches(w)
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = C_DARK
            p.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_BG
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = C_BODY
    return table_shape


# ═══════════════════════════════════════════
#  SLIDE 1: COVER
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_DARK)

# Gradient-like accent bar
add_rect(s, 0, 0, W.inches, 0.08, C_BRAND, None)

# Icon circle
icon_shape = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.1), Inches(1.6), Inches(1.2), Inches(1.2))
icon_shape.fill.solid()
icon_shape.fill.fore_color.rgb = C_BRAND
icon_shape.line.fill.background()
tf = icon_shape.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.text = "⚙"
p.font.size = Pt(42)
p.font.color.rgb = C_WHITE
p.alignment = PP_ALIGN.CENTER

add_text_box(s, 0.5, 3.1, 12.3, 1.2, "iHealthSim", font_size=48,
             color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, 0.5, 4.3, 12.3, 0.8,
             "工业设备健康状态评估系统\n全仿真 · 完整链路 · 实时看板",
             font_size=20, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 0.5, 5.8, 12.3, 0.5, "版本 0.1.0  ·  2025-05",
             font_size=13, color=C_MUTED, align=PP_ALIGN.CENTER)

# ── Slide number indicator ──
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "1 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 2: SECTION - Project Overview
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_DARK)
add_text_box(s, 1, 2.6, 11.3, 1.2, "项目概述", font_size=42,
             color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.9, 11.3, 0.8,
             "基于仿真设备 + MQTT + 决策树 + Vue 前端的\n工业设备健康评估原型系统",
             font_size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "2 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 3: Core Values
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_BG)
add_text_box(s, 0.8, 0.4, 11, 0.6, "项目定位与核心价值", font_size=28,
             color=C_DARK, bold=True)
add_text_box(s, 0.8, 0.9, 11, 0.4, "面向 IIoT / PdM / EHM 的完整技术链路验证原型",
             font_size=13, color=C_MUTED)

# 4 cards
add_card(s, 0.8, 1.5, 5.8, 2.2,
         "零硬件依赖",
         ["所有设备数据由物理仿真引擎生成",
          "覆盖 6 种工况循环、渐进退化",
          "故障注入、瞬态尖峰等真实场景",
          "无需连接真实 PLC/传感器"],
         icon_color=C_BRAND, title_size=15, body_size=11)

add_card(s, 7.0, 1.5, 5.8, 2.2,
         "完整数据链路",
         ["仿真 → MQTT 传输 → 采集落盘",
          "→ 特征工程 → 决策树训练",
          "→ 在线打分 → SSE 推送",
          "→ Web 实时看板展示"],
         icon_color=C_GREEN, title_size=15, body_size=11)

add_card(s, 0.8, 4.0, 5.8, 2.2,
         "可审计预测",
         ["决策树输出完整决策路径",
          "从根节点到叶节点逐步解释",
          "前端渲染为诊断依据时间线",
          "非黑盒模型，安全可信"],
         icon_color=C_AMBER, title_size=15, body_size=11)

add_card(s, 7.0, 4.0, 5.8, 2.2,
         "工业级可扩展",
         ["标准 MQTT 消息格式",
          "无缝替换为真实设备数据源",
          "Pipeline 架构支持 sklearn",
          "模型热替换，即插即用"],
         icon_color=C_PURPLE, title_size=15, body_size=11)

# Bottom tag
add_rect_text(s, 0.8, 6.5, 2.6, 0.35, "IIoT 原型验证", C_LIGHT_BLUE, RGBColor(0x93,0xC5,0xFD), 10, C_BRAND, True)
add_rect_text(s, 3.6, 6.5, 2.6, 0.35, "预测性维护 PdM", C_LIGHT_GREEN, RGBColor(0x6E,0xE7,0xB7), 10, RGBColor(0x06,0x5F,0x46), True)
add_rect_text(s, 6.4, 6.5, 2.6, 0.35, "设备健康管理 EHM", C_LIGHT_AMBER, RGBColor(0xFC,0xD3,0x4D), 10, RGBColor(0x92,0x40,0x0E), True)

add_text_box(s, 11.5, 7.1, 1.5, 0.3, "3 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 4: SECTION - Architecture
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_DARK)
add_text_box(s, 1, 2.6, 11.3, 1.2, "系统架构", font_size=42,
             color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.9, 11.3, 0.8,
             "三层架构：仿真设备层 → MQTT 传输层 → 计算 + 展示层",
             font_size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "4 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 5: Architecture Diagram
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_BG)
add_text_box(s, 0.8, 0.4, 11, 0.6, "系统架构全景图", font_size=28, color=C_DARK, bold=True)
add_text_box(s, 0.8, 0.9, 11, 0.3, "三层解耦：设备仿真 → MQTT 消息中间件 → 后端计算 + 前端展示", font_size=12, color=C_MUTED)

# Frontend
add_rect_text(s, 4.0, 1.3, 5.3, 0.55, "Vue 3 前端 :5173   Login / Dashboard / Admin",
              C_LIGHT_PURPLE, RGBColor(0xC4,0xB5,0xFD), 13, C_PURPLE, True)
# Arrow down
add_text_box(s, 6.4, 1.85, 0.5, 0.35, "▼", font_size=16, color=C_MUTED, align=PP_ALIGN.CENTER)

# Backend group
add_rect(s, 1.5, 2.15, 10.3, 1.25, C_WHITE, RGBColor(0xBF, 0xDB, 0xFE), Pt(2))
add_text_box(s, 1.7, 2.15, 3, 0.3, "Flask 后端 :5000", font_size=11, color=C_BRAND, bold=True)
for i, (label, clr, bdr) in enumerate([
    ("REST API", C_LIGHT_BLUE, RGBColor(0x93,0xC5,0xFD)),
    ("SSE Hub", C_LIGHT_PURPLE, RGBColor(0xC4,0xB5,0xFD)),
    ("MQTT Sub", C_LIGHT_GREEN, RGBColor(0x6E,0xE7,0xB7)),
    ("State", C_LIGHT_AMBER, RGBColor(0xFC,0xD3,0x4D)),
    ("Scorer", C_LIGHT_RED, RGBColor(0xFC,0xA5,0xA5)),
    ("Auth", RGBColor(0xFC,0xE7,0xF3), RGBColor(0xF9,0xA8,0xD4)),
]):
    add_rect_text(s, 2.2 + i * 1.5, 2.55, 1.3, 0.65, label, clr, bdr, 11, C_BODY, True)

# Arrow
add_text_box(s, 6.4, 3.4, 0.5, 0.3, "▼", font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 2.5, 3.4, 2, 0.3, "subscribe", font_size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)

# EMQX
add_rect_text(s, 3.5, 3.65, 6.3, 0.6, "EMQX Broker :1883   topic: telemetry/raw/#",
              C_LIGHT_GREEN, RGBColor(0x6E,0xE7,0xB7), 13, RGBColor(0x06,0x5F,0x46), True)

# Arrow
add_text_box(s, 6.4, 4.25, 0.5, 0.3, "▼", font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 2.5, 4.25, 2, 0.3, "publish", font_size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)

# Simulator
add_rect_text(s, 4.0, 4.5, 5.3, 0.6, "仿真设备 × 3    PUMP-001  ~  PUMP-003",
              C_LIGHT_BLUE, RGBColor(0x93,0xC5,0xFD), 13, C_BRAND, True)

# Insight box
add_rect(s, 0.8, 5.5, 11.7, 1.2, RGBColor(0xFF, 0xFB, 0xEB), RGBColor(0xFD, 0xE6, 0x8A), Pt(1))
tf = add_rect(s, 0.8, 5.5, 0.08, 1.2, C_AMBER, None).text_frame  # left accent bar
tx = add_text_box(s, 1.1, 5.55, 11.2, 1.1,
                  "架构原则\n设备不感知消费者，后端不感知发布者。通过 EMQX 解耦，单台设备可被 N 个后端同时消费；"
                  "替换真实设备时只需发布相同格式的 JSON 到相同 topic。",
                  font_size=11, color=RGBColor(0x92,0x40,0x0E))
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "5 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 6: Data Flow
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_BG)
add_text_box(s, 0.8, 0.4, 11, 0.6, "数据流详解", font_size=28, color=C_DARK, bold=True)
add_text_box(s, 0.8, 0.9, 11, 0.3, "从仿真器 step() 一个方法调用到前端看板实时刷新", font_size=12, color=C_MUTED)

# Left: flow steps
steps = [
    "1   DeviceSimulator.step() 更新健康度、计算 6 个遥测点位",
    "2   封装 TelemetryPoint → JSON → MQTT publish",
    "3   EMQX Broker 路由到所有 subscriber",
    "4   MqttSubscriber 解析 JSON → InMemoryState.ingest()",
    "5   SSEHub.push(\"telemetry\") → 前端实时刷新",
    "6   ScorerWorker 滑动窗口触发 pipeline.predict()",
    "7   去抖确认 → HealthResult → SSEHub.push(\"prediction\")",
]
y = 1.5
for step in steps:
    add_text_box(s, 0.8, y, 6, 0.35, step, font_size=12, color=C_BODY)
    y += 0.52

# Right: key points
add_card(s, 7.3, 1.5, 5.3, 2.0, "关键技术点", [
    "• 以 temp_c 点位作为窗口触发器，每窗口只打分一次",
    "• 去抖逻辑：恶化连续 3 次确认，恢复连续 5 次确认",
    "• SSE 基于 HTTP，浏览器原生支持断线重连",
    "• Thread-safe：ScorerWorker 使用锁保护模型推理",
], title_size=14, body_size=11)

# Right bottom: code block
add_code_block(s, 7.3, 3.8, 5.3, 2.2,
    '# MQTT 消息格式（JSON）\n'
    '{\n'
    '  "ts_ms": 1710000000000,\n'
    '  "asset_id": "PUMP-001",\n'
    '  "point": "temp_c",\n'
    '  "value": 55.2,\n'
    '  "quality": "good"\n'
    '}')
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "6 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 7: SECTION - Core Modules
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_DARK)
add_text_box(s, 1, 2.6, 11.3, 1.2, "核心模块详解", font_size=42,
             color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.9, 11.3, 0.8,
             "仿真引擎 · 特征工程 · 决策树 · 在线打分 · 去抖迟滞",
             font_size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "7 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 8: Simulator + Features
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_BG)
add_text_box(s, 0.8, 0.4, 11, 0.6, "设备仿真器 + 特征工程", font_size=28, color=C_DARK, bold=True)
add_text_box(s, 0.8, 0.9, 11, 0.3, "物理驱动的渐进退化仿真模型", font_size=12, color=C_MUTED)

# Simulation model
add_card(s, 0.8, 1.4, 5.8, 2.3, "仿真模型 (DeviceSimulator)", [
    "• latent_health [0→1] 潜在健康度变量",
    "• 基础退化: degradation/hr + 高斯噪声",
    "• 故障注入: 3× 加速退化",
    "• 随机恢复: 5%/s 概率微恢复",
    "• 6 工况循环: 1500-2900 RPM",
    "• 环境温度正弦漂移 ±3°C",
    "• 传感器噪声随恶化 ×3 扩大",
    "• 瞬态尖峰 0.05%/s 概率",
], title_size=14, body_size=10)

# Feature engineering
add_card(s, 7.0, 1.4, 5.8, 2.3, "特征工程流水线", [
    "① long → wide pivot",
    "② 60s 滑动窗口划分",
    "③ 基础统计: mean/std/min/max/p95",
    "④ 趋势特征: 线性回归斜率",
    "⑤ Δ 特征: 与上一窗口差值",
    "⑥ 工况归一化振动: vib/(rpm/1000)",
    "⑦ 标签: 窗口内最大 health_level",
    "输出: data/features.csv (~30 列特征 + y)",
], title_size=14, body_size=10)

# Telemetry points table
add_text_box(s, 0.8, 3.9, 5, 0.3, "遥测点位", font_size=14, color=C_DARK, bold=True)
add_table(s, 0.8, 4.2, [2.0, 2.0, 1.8, 1.8],
          ["点位", "含义", "正常范围", "恶化趋势"],
          [["rpm", "转速", "1500-2900", "基准恒定"],
           ["load", "负载比", "0.1-1.0", "基准恒定"],
           ["vib_rms", "振动 RMS", "~1.2-3.5", "↑ 显著上升"],
           ["temp_c", "温度 (°C)", "~35-65", "↑ 上升"],
           ["motor_current_a", "电机电流 (A)", "~10-40", "↑ 上升"],
           ["label_health_level", "健康标签", "0-3", "训练用"]])

# Health levels
add_text_box(s, 7.0, 3.9, 5, 0.3, "健康等级定义", font_size=14, color=C_DARK, bold=True)
add_table(s, 7.0, 4.2, [1.4, 1.8, 1.8],
          ["等级", "latent_health", "状态"],
          [["Lv0 健康", "≥ 0.80", "🟢 正常运转"],
           ["Lv1 注意", "0.60~0.80", "🟡 关注趋势"],
           ["Lv2 警告", "0.40~0.60", "🔴 明显异常"],
           ["Lv3 危险", "< 0.40", "🟣 停机检查"]])

add_text_box(s, 11.5, 7.1, 1.5, 0.3, "8 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 9: Training + Online Scoring
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_BG)
add_text_box(s, 0.8, 0.4, 11, 0.6, "模型训练 + 在线打分", font_size=28, color=C_DARK, bold=True)
add_text_box(s, 0.8, 0.9, 11, 0.3, "scikit-learn Pipeline + 去抖迟滞逻辑", font_size=12, color=C_MUTED)

# Left: Training
add_card(s, 0.8, 1.4, 5.8, 1.5, "训练 Pipeline", [
    "ColumnTransformer(SimpleImputer(median))",
    "        ↓",
    "DecisionTreeClassifier(max_depth=5, min_leaf=50, balanced)",
    "时间序列分割: 按窗口时间顺序切分，避免数据泄露",
    "feature_set=online: 只用 OnlineScorer 可实时计算的特征",
], title_size=14, body_size=11)

# 8 online features
add_card(s, 0.8, 3.1, 5.8, 1.2, "8 个在线特征", [
    "rpm_mean  |  load_mean  |  vib_rms_mean  |  temp_c_mean",
    "motor_current_a_mean  |  vib_rms_std  |  temp_c_std",
    "vib_rms_norm (= vib_rms / (rpm/1000))",
], title_size=14, body_size=11)

# Online scorer flow
add_card(s, 7.0, 1.4, 5.8, 2.0, "OnlineScorer 工作流程", [
    "① 按 asset_id 缓冲遥测点队列",
    "② temp_c 触发 60s 窗口边界判定",
    "③ 提取窗口数据 → 计算 8 个实时特征",
    "④ pipeline.predict() + predict_proba()",
    "⑤ 去抖确认 → HealthResult",
    "⑥ 遍历决策树生成可读解释路径",
], title_size=14, body_size=11)

# Debounce config
add_rect(s, 7.0, 3.6, 5.8, 2.5, C_LIGHT_AMBER, RGBColor(0xFD, 0xE6, 0x8A), Pt(1.5))
add_text_box(s, 7.2, 3.65, 5.4, 0.3, "去抖配置 (DebounceConfig)", font_size=13, color=RGBColor(0x92,0x40,0x0E), bold=True)
add_table(s, 7.2, 4.0, [1.5, 0.8, 3.0],
          ["参数", "值", "含义"],
          [["raise_n", "3", "连续 3 次确认等级上升（快速报警）"],
           ["recover_n", "5", "连续 5 次确认等级恢复（谨慎解除）"],
           ["abnormal_level", "2", "等级 ≥2 被视为异常"]])

add_rect(s, 7.2, 5.2, 5.4, 0.7, RGBColor(0xFF, 0xFB, 0xEB), None)
add_text_box(s, 7.4, 5.3, 5.0, 0.5, "工业语义：报警需快速响应，解除需谨慎验证 —— 宁可多等，不误解除",
             font_size=11, color=RGBColor(0x92,0x40,0x0E))

add_text_box(s, 11.5, 7.1, 1.5, 0.3, "9 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 10: SECTION - Backend + Frontend
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_DARK)
add_text_box(s, 1, 2.6, 11.3, 1.2, "后端与前端的协作", font_size=42,
             color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.9, 11.3, 0.8,
             "Flask + SSE + Vue 3 实时看板",
             font_size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "10 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 11: Backend + Frontend Detail
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_BG)
add_text_box(s, 0.8, 0.4, 11, 0.6, "后端模块 + 前端看板", font_size=28, color=C_DARK, bold=True)
add_text_box(s, 0.8, 0.9, 11, 0.3, "7 个后端模块 · 4 个前端页面 · 实时事件驱动", font_size=12, color=C_MUTED)

# Backend modules
for i, (title, body) in enumerate([
    ("Flask 主应用 (app.py)", "启动时自动初始化 MQTT 订阅器、状态管理、SSE 中心、打分器；并拉起 3 台仿真设备后台进程"),
    ("认证授权 (auth.py)", "JWT HS256 认证 + 角色管理 (admin/operator) + 设备级 ACL 权限控制，默认 admin/admin123"),
    ("SSE 事件中心 (sse.py)", "发布-订阅模式，多客户端并发，自动丢弃慢消费者过期消息，3 类事件: telemetry/prediction/flow"),
    ("MQTT 订阅器 + State", "paho-mqtt 异步 loop，解析 JSON → TelemetryPoint → State + SSE + Scorer 三路分发"),
]):
    cx = 0.8 if i % 2 == 0 else 7.0
    cy = 1.4 + (i // 2) * 1.3
    add_card(s, cx, cy, 5.8, 1.1, title, [body], title_size=13, body_size=10)

# Frontend routes
add_text_box(s, 0.8, 4.2, 5, 0.3, "前端路由设计", font_size=14, color=C_DARK, bold=True)
add_table(s, 0.8, 4.5, [1.2, 1.8, 1.2],
          ["路径", "组件", "认证"],
          [["/login", "Login.vue", "Guest"],
           ["/register", "Register.vue", "Guest"],
           ["/dashboard", "Dashboard.vue", "Token"],
           ["/admin", "Admin.vue", "Admin"]])

# Dashboard features
add_card(s, 7.0, 4.2, 5.8, 2.5, "Dashboard 核心页面功能", [
    "• 状态指示条：MQTT / 模型 / 数据状态",
    "• 设备标签栏：多设备切换 + 健康等级徽章",
    "• 健康分数环形图 + 四级概率分布柱状图",
    "• 诊断依据时间线（决策路径可视化）",
    "• 实时遥测表 + 事件日志",
    "• 诊断报告：设备概览/异常指标/维护建议",
    "• 「恶化重演」：一键重启完整退化过程",
], title_size=13, body_size=10)

add_text_box(s, 11.5, 7.1, 1.5, 0.3, "11 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 12: SECTION - Design Decisions
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_DARK)
add_text_box(s, 1, 2.6, 11.3, 1.2, "关键技术决策", font_size=42,
             color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.9, 11.3, 0.8,
             "为什么选择这些技术栈？每个选择背后都有明确的工业场景考量",
             font_size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "12 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 13: Design Decisions Detail
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_BG)
add_text_box(s, 0.8, 0.4, 11, 0.6, "设计决策与理由", font_size=28, color=C_DARK, bold=True)
add_text_box(s, 0.8, 0.9, 11, 0.3, "每个技术选择背后都有明确的工业场景考量", font_size=12, color=C_MUTED)

decisions = [
    (C_BRAND, "为什么用决策树而非深度学习？",
     "可解释性强：输出完整决策路径，非黑盒；数据量（千-万级窗口）决策树足够；单个 joblib 文件部署，无需 GPU，推理毫秒级"),
    (C_GREEN, "为什么用 MQTT 而非 HTTP？",
     "IIoT 事实标准：发布/订阅解耦设备与消费者；QoS 消息保证；二进制协议带宽低；与 EMQX/HiveMQ/AWS IoT 生态兼容"),
    (C_AMBER, "为什么用 SSE 而非 WebSocket？",
     "前端只需单向接收推送；SSE 基于 HTTP 无需协议升级；浏览器原生断线重连；实现复杂度远低于 WebSocket"),
    (C_RED, "为什么需要去抖迟滞？",
     "传感器噪声和瞬态工况会导致等级边界抖动；上升快速确认（3次），恢复慢速确认（5次），符合工业\"报警易·解除难\"需求"),
    (C_PURPLE, "为什么用时间序列分割？",
     "避免随机切分导致未来信息泄露到训练集，使评估结果更接近真实在线部署场景，每个 asset 独立切分保证评估严谨性"),
    (RGBColor(0xEC, 0x48, 0x99), "为什么 feature_set=online？",
     "训练时只使用在线可实时计算的特征列，避免大量离线特征被中位数填充成常数，导致预测结果长期不变化"),
]

for i, (accent, title, body) in enumerate(decisions):
    cx = 0.8 if i % 2 == 0 else 7.0
    cy = 1.4 + (i // 2) * 1.8
    add_rect(s, cx, cy, 5.8, 1.6, C_WHITE, C_BORDER)
    # accent top border
    add_rect(s, cx, cy, 5.8, 0.05, accent, None)
    add_text_box(s, cx + 0.2, cy + 0.12, 5.4, 0.3, title, font_size=13, color=C_DARK, bold=True)
    add_text_box(s, cx + 0.2, cy + 0.55, 5.4, 0.95, body, font_size=11, color=C_MUTED)

# Insight
add_rect(s, 0.8, 6.9, 11.7, 0.0, C_AMBER, None)
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "13 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 14: SECTION - Deploy & Extend
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_DARK)
add_text_box(s, 1, 2.6, 11.3, 1.2, "部署与扩展", font_size=42,
             color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.9, 11.3, 0.8,
             "一键启动 · Docker 部署 · 接入真实设备 · 模型替换",
             font_size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "14 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 15: Deploy & Extend Detail
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_BG)
add_text_box(s, 0.8, 0.4, 11, 0.6, "部署方案 + 扩展指南", font_size=28, color=C_DARK, bold=True)
add_text_box(s, 0.8, 0.9, 11, 0.3, "从开发到生产，从仿真到真实", font_size=12, color=C_MUTED)

# One-click start
add_code_block(s, 0.8, 1.4, 5.8, 0.7,
    '# 一键启动\n'
    'bash start.sh\n'
    '# 前端 http://localhost:5173\n'
    '# 后端 http://localhost:5000\n'
    '# 账号 admin / admin123')

# Docker compose
add_code_block(s, 7.0, 1.4, 5.8, 2.2,
    '# Docker Compose\n'
    'services:\n'
    '  emqx:   image: emqx/emqx\n'
    '  mysql:  image: mysql:8\n'
    '          environment:\n'
    '            MYSQL_ROOT_PASSWORD: root\n'
    '            MYSQL_DATABASE: ihealthsim\n'
    '  backend:  build: .\n'
    '  frontend: build: ./frontend')

# Real device integration
add_card(s, 0.8, 3.8, 5.8, 1.8, "接入真实设备 — 零代码改动", [
    "只需确保 MQTT 消息格式一致：",
    '{',
    '  "ts_ms": 1710000000000,',
    '  "asset_id": "REAL-PUMP-001",',
    '  "point": "temp_c", "value": 55.2,',
    '  "quality": "good"',
    '}',
    "系统自动接收 → 存储 → 打分 → 推送",
], title_size=14, body_size=9)

# Extend
add_card(s, 7.0, 3.8, 5.8, 1.8, "扩展方向", [
    "🔧 替换模型：任何 sklearn 兼容分类器",
    "  → 修改 train.py + 对齐 ONLINE_FEATURE_COLS",
    "",
    "📊 增加点位：四点对齐",
    "  → step() → features.py → train.py → Dashboard.vue",
    "",
    "📈 未来规划",
    "  → XGBoost/LightGBM · ECharts 历史趋势",
    "  → 报警规则引擎 · Docker 一键部署 · i18n",
], title_size=14, body_size=9)

add_text_box(s, 11.5, 7.1, 1.5, 0.3, "15 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 16: Tech Stack Summary
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_BG)
add_text_box(s, 0.8, 0.4, 11, 0.6, "技术栈总览", font_size=28, color=C_DARK, bold=True)
add_text_box(s, 0.8, 0.9, 11, 0.3, "项目依赖全景图", font_size=12, color=C_MUTED)

# 3 columns
stacks = [
    ("Python 核心", C_BRAND, ["numpy + pandas 数据处理", "scikit-learn 决策树", "joblib 模型持久化", "rich CLI 美化"]),
    ("网络与通信", C_GREEN, ["EMQX MQTT Broker", "paho-mqtt 客户端", "Flask REST API", "SSE 事件推送"]),
    ("前端 + 数据库", C_PURPLE, ["Vue 3 + Vue Router", "Vite 构建工具", "MySQL 用户/事件存储", "JWT (PyJWT) 认证"]),
]
for i, (title, color, items) in enumerate(stacks):
    cx = 0.8 + i * 4.2
    add_card(s, cx, 1.4, 3.9, 2.2, title, ["• " + it for it in items], title_size=14, body_size=11)

# Project scale
add_text_box(s, 0.8, 3.9, 5, 0.3, "项目规模", font_size=14, color=C_DARK, bold=True)
metrics = [
    ("12+", "CLI 子命令", C_BRAND),
    ("10+", "REST API 端点", C_GREEN),
    ("7", "后端模块", C_PURPLE),
    ("4", "前端页面", C_AMBER),
]
for i, (num, label, color) in enumerate(metrics):
    cx = 0.8 + i * 3.1
    shape = add_rect(s, cx, 4.3, 2.8, 1.2, C_WHITE, C_BORDER)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(32)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = C_MUTED
    p2.alignment = PP_ALIGN.CENTER

# Key libs
add_text_box(s, 0.8, 5.8, 5, 0.3, "关键依赖", font_size=14, color=C_DARK, bold=True)
add_table(s, 0.8, 6.1, [3.0, 3.0, 3.0, 3.0],
          ["Python", "版本", "Node.js", "版本"],
          [["numpy", "≥2.0", "vue", "3.x"],
           ["pandas", "≥2.2", "vue-router", "4.x"],
           ["scikit-learn", "≥1.5", "vite", "5.x"],
           ["flask", "≥3.0", "", ""],
           ["paho-mqtt", "≥2.1", "", ""]])

add_text_box(s, 11.5, 7.1, 1.5, 0.3, "16 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SLIDE 17: END
# ═══════════════════════════════════════════
s = add_blank_slide()
fill_bg(s, C_DARK)
add_rect(s, 0, 0, W.inches, 0.08, C_BRAND, None)

icon_shape = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.8), Inches(1.5), Inches(1.5))
icon_shape.fill.solid()
icon_shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
icon_shape.line.color.rgb = RGBColor(0x33, 0x44, 0x55)
icon_shape.line.width = Pt(1)
tf = icon_shape.text_frame
p = tf.paragraphs[0]
p.text = "⚙"
p.font.size = Pt(48)
p.font.color.rgb = C_BRAND
p.alignment = PP_ALIGN.CENTER

add_text_box(s, 0.5, 3.6, 12.3, 1.2, "谢谢", font_size=48,
             color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(s, 0.5, 4.6, 12.3, 0.8,
             "iHealthSim — 工业设备健康状态评估系统\n从仿真到实时看板的完整链路",
             font_size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 0.5, 5.8, 12.3, 0.5,
             "详细文档: doc/ARCHITECTURE.md  |  一键启动: bash start.sh",
             font_size=13, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text_box(s, 11.5, 7.1, 1.5, 0.3, "17 / 17", font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════
output_path = "/Users/huyuuu/ Equipment Health Assessment System/doc/iHealthSim.pptx"
prs.save(output_path)
print(f"PPTX saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
